import unittest
import os
import io
from src.document_processor import DocumentProcessor
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

class TestDocumentProcessor(unittest.TestCase):
    def setUp(self):
        self.processor = DocumentProcessor()

    def create_mock_file(self, content, filename):
        file = io.BytesIO(content)
        file.name = filename
        file.size = len(content)
        return file
    
    def create_txt_file(self, text, filename="test.txt"):
        return self.create_mock_file(text.encode("utf-8"), filename)

    def create_docx_file(self, text, filename="test.docx"):
        buffer = io.BytesIO()
        doc = DocxDocument()
        doc.add_paragraph(text)
        doc.save(buffer)
        buffer.seek(0)
        buffer.name = filename
        buffer.size = buffer.getbuffer().nbytes
        return buffer

    def create_pptx_file(self, text, filename="test.pptx"):
        buffer = io.BytesIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = text
        prs.save(buffer)
        buffer.seek(0)
        buffer.name = filename
        buffer.size = buffer.getbuffer().nbytes
        return buffer

    def create_excel_file(self, text, filename="test.xlsx"):
        buffer = io.BytesIO()
        df = pd.DataFrame({'Col1': [text]})
        with pd.ExcelWriter(buffer, engine='xlsxwriter') as writer:
            df.to_excel(writer, index=False)
        buffer.seek(0)
        buffer.name = filename
        buffer.size = buffer.getbuffer().nbytes
        return buffer

    def test_process_txt(self):
        file = self.create_txt_file("Hello World")
        docs = self.processor.process_files([file])
        self.assertEqual(len(docs), 1)
        self.assertIn("Hello World", docs[0].page_content)
        self.assertEqual(docs[0].metadata['type'], 'txt')

    def test_process_docx(self):
        file = self.create_docx_file("Hello Docx")
        docs = self.processor.process_files([file])
        self.assertEqual(len(docs), 1)
        self.assertIn("Hello Docx", docs[0].page_content)
        self.assertEqual(docs[0].metadata['type'], 'docx')
        
    def test_process_pptx(self):
        # python-pptx usually needs a real file or correct BytesIO handling.
        # Our implementation uses prs = Presentation(file) which accepts file object.
        file = self.create_pptx_file("Hello PPTX")
        docs = self.processor.process_files([file])
        self.assertEqual(len(docs), 1)
        self.assertIn("Hello PPTX", docs[0].page_content)
        self.assertEqual(docs[0].metadata['type'], 'pptx')

    def test_process_excel(self):
        file = self.create_excel_file("Hello Excel")
        docs = self.processor.process_files([file])
        self.assertEqual(len(docs), 1)
        # Excel content is markdown table
        self.assertIn("Hello Excel", docs[0].page_content)
        self.assertEqual(docs[0].metadata['type'], 'excel')

if __name__ == '__main__':
    unittest.main()
