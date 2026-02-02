import fitz  # PyMuPDF
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter

from langchain.docstore.document import Document
from typing import List, Dict, Any
from src.config import AppConfig
import io
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation

class ProcessingLogger:
    """Tracks document processing stats for user visibility."""
    
    def __init__(self):
        self.logs: List[Dict[str, Any]] = []

    def log(self, file_name: str, step: str, details: str):
        entry = {
            "file": file_name,
            "step": step,
            "details": details
        }
        self.logs.append(entry)
        # In a real app, you might also push this to a UI stream or database
        print(f"[{file_name}] {step}: {details}")

class DocumentProcessor:
    """Handles document ingestion and processing for multiple formats."""

    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=AppConfig.CHUNK_SIZE,
            chunk_overlap=AppConfig.CHUNK_OVERLAP
        )
        self.logger = ProcessingLogger()

    def process_files(self, uploaded_files) -> List[Document]:
        """
        Main dispatcher for processing uploaded files.
        """
        all_documents = []
        
        for file in uploaded_files:
            file_name = file.name
            file_ext = file_name.split('.')[-1].lower()
            self.logger.log(file_name, "Ingestion", f"Started processing. Size: {file.size / 1024:.2f} KB")

            try:
                docs = []
                if file_ext == 'pdf':
                    docs = self._process_pdf(file)
                elif file_ext in ['docx', 'doc']:
                    docs = self._process_docx(file)
                elif file_ext in ['pptx', 'ppt']:
                    docs = self._process_pptx(file)
                elif file_ext in ['xlsx', 'xls']:
                    docs = self._process_excel(file)
                elif file_ext == 'txt':
                    docs = self._process_txt(file)
                else:
                    self.logger.log(file_name, "Error", f"Unsupported format: {file_ext}")
                    continue
                
                if docs:
                    self.logger.log(file_name, "Extraction", f"Extracted {len(docs)} pages/sections.")
                    all_documents.extend(docs)
                else:
                    self.logger.log(file_name, "Warning", "No text extracted.")

            except Exception as e:
                self.logger.log(file_name, "Error", f"Processing failed: {str(e)}")
        
        return all_documents

    def _process_pdf(self, file) -> List[Document]:
        file.seek(0)
        pdf_bytes = file.read()
        documents = []
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={"source": file.name, "page": page_num + 1, "type": "pdf"}
                    ))
        return documents

    def _process_docx(self, file) -> List[Document]:
        file.seek(0)
        doc = DocxDocument(file)
        text_content = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        # Docx is usually treated as a single continuous doc, but we can try to split by logic if needed.
        # For now, we return as one document per file, or we could split by headers.
        full_text = "\n".join(text_content)
        return [Document(
            page_content=full_text,
            metadata={"source": file.name, "type": "docx"}
        )] if full_text.strip() else []

    def _process_pptx(self, file) -> List[Document]:
        file.seek(0)
        prs = Presentation(file)
        documents = []
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            
            slide_text = "\n".join(text_runs)
            if slide_text.strip():
                documents.append(Document(
                    page_content=slide_text,
                    metadata={"source": file.name, "page": i + 1, "type": "pptx"}
                ))
        return documents

    def _process_excel(self, file) -> List[Document]:
        file.seek(0)
        # Read all sheets
        xls = pd.ExcelFile(file, engine='openpyxl')
        documents = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            # Convert to markdown or text representation
            text_content = f"Sheet: {sheet_name}\n"
            text_content += df.to_markdown(index=False)
            
            documents.append(Document(
                page_content=text_content,
                metadata={"source": file.name, "sheet": sheet_name, "type": "excel"}
            ))
        return documents

    def _process_txt(self, file) -> List[Document]:
        file.seek(0)
        text = file.read().decode("utf-8")
        return [Document(
            page_content=text,
            metadata={"source": file.name, "type": "txt"}
        )]

    def chunk_documents(self, documents: List[Document]) -> List[Document]:
        """
        Splits documents into smaller chunks.
        """
        chunks = self.text_splitter.split_documents(documents)
        
        # Log chunking stats if possible (would need refactoring to pass logger here or make logger global/singleton)
        # For now, we assume the caller will log the result size.
        return chunks
